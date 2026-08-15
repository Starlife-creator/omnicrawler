"""document_ir 富文档解析（S2）：.docx / .pptx / .odt / .epub。

依赖策略：
- .docx → python-docx（懒加载，缺失时抛 ModuleNotFoundError 提示 omnicrawl[document]）
- .pptx → python-pptx（懒加载，同上）
- .odt  → zipfile + xml.etree.ElementTree（标准库，零依赖）
- .epub → zipfile + 项目内 html_tools（标准库，零依赖）

注册方式与 parsers.py 一致：register_document_parser 装饰器。
"""

from __future__ import annotations

import io
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path
from typing import Any

from ..extraction.html_tools import node_text, parse_html, select_nodes
from .base import DocumentIR
from .parsers import _require_file, _require_parser, register_document_parser

#: 正文候选容器（与 parsers.py 的 HTML 解析一致）
_CONTENT_SELECTOR = "p,li,blockquote,h1,h2,h3,h4,h5,h6"

#: ODT 文本命名空间（段落标签）
_ODT_TEXT_NS = "{urn:oasis:names:tc:opendocument:xmlns:text:1.0}"


def _dedupe(items: list[str]) -> list[str]:
    """顺序去重（HTML 各节点文本可能重叠）。"""
    out: list[str] = []
    for item in items:
        if item and item not in out:
            out.append(item)
    return out


# ── .docx ────────────────────────────────────────────────
@register_document_parser(".docx")
def _parse_docx(path: Path, options: dict[str, Any]) -> DocumentIR:
    _require_file(path)
    try:
        import docx  # python-docx
    except ImportError:
        _require_parser(".docx", "python-docx")

    document = docx.Document(str(path))
    paragraphs = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    tables = [
        [[cell.text.strip() for cell in row.cells] for row in table.rows]
        for table in document.tables
    ]

    core = document.core_properties
    title = paragraphs[0] if paragraphs else (core.title or path.stem)
    if title == paragraphs[0] and paragraphs:
        paragraphs = paragraphs[1:]
    elif not paragraphs:
        title = path.stem

    metadata: dict[str, Any] = {"author": core.author or ""}
    if core.created is not None:
        metadata["created"] = core.created.isoformat()
    if core.modified is not None:
        metadata["modified"] = core.modified.isoformat()
    metadata["paragraph_count"] = len(paragraphs)

    return DocumentIR(
        source=path,
        kind=".docx",
        title=title,
        paragraphs=paragraphs,
        tables=tables,
        metadata=metadata,
    )


# ── .pptx ────────────────────────────────────────────────
@register_document_parser(".pptx")
def _parse_pptx(path: Path, options: dict[str, Any]) -> DocumentIR:
    _require_file(path)
    try:
        from pptx import Presentation
    except ImportError:
        _require_parser(".pptx", "python-pptx")

    prs = Presentation(str(path))
    paragraphs: list[str] = []
    tables: list[list[list[str]]] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        paragraphs.append(text)
            elif shape.has_table:
                tables.append(
                    [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                )

    title = paragraphs[0] if paragraphs else path.stem
    if title == paragraphs[0] and paragraphs:
        paragraphs = paragraphs[1:]

    return DocumentIR(
        source=path,
        kind=".pptx",
        title=title,
        paragraphs=paragraphs,
        tables=tables,
        metadata={"slide_count": len(prs.slides), "paragraph_count": len(paragraphs)},
    )


# ── .odt ─────────────────────────────────────────────────
@register_document_parser(".odt")
def _parse_odt(path: Path, options: dict[str, Any]) -> DocumentIR:
    _require_file(path)
    try:
        with zipfile.ZipFile(path) as zf:
            content = _read_zip_entry(zf, "content.xml")
    except (zipfile.BadZipFile, KeyError) as exc:
        raise ValueError(f"document_ir: 无效的 ODT 文件: {path} ({exc})") from exc

    tree = ET.parse(io.BytesIO(content))
    paragraphs: list[str] = []
    for elem in tree.iter(_ODT_TEXT_NS + "p"):
        text = "".join(elem.itertext()).strip()
        if text:
            paragraphs.append(" ".join(text.split()))

    title = paragraphs[0] if paragraphs else path.stem
    if title == paragraphs[0] and paragraphs:
        paragraphs = paragraphs[1:]

    return DocumentIR(
        source=path,
        kind=".odt",
        title=title,
        paragraphs=paragraphs,
        metadata={"paragraph_count": len(paragraphs)},
    )


# ── .epub ────────────────────────────────────────────────
# B07-002：zip 单条目解压上限（64 MiB）——防 zip 炸弹/超大 content.xml 内存耗尽
_ZIP_ENTRY_MAX_BYTES = 64 * 1024 * 1024


def _read_zip_entry(
    zf: zipfile.ZipFile, name: str, *, max_bytes: int = _ZIP_ENTRY_MAX_BYTES,
) -> bytes:
    """解压 zip 条目前检查 file_size 上限（defense-in-depth）。"""
    info = zf.getinfo(name)
    if info.file_size > max_bytes:
        raise ValueError(
            f"document_ir: zip 条目过大（{info.file_size} 字节 > 上限 {max_bytes}）: {name}"
        )
    return zf.read(name)


def _epub_spine_order(zf: zipfile.ZipFile) -> list[str]:
    """按 EPUB 阅读顺序返回内容文件相对路径列表；无 OPF 时返回空。"""
    try:
        container = ET.parse(io.BytesIO(_read_zip_entry(zf, "META-INF/container.xml")))
    except (KeyError, ET.ParseError):
        return []
    ns = {"c": "urn:oasis:names:tc:opendocument:xmlns:container"}
    rootfile = container.find(".//c:rootfile", ns)
    if rootfile is None:
        return []
    opf_path = rootfile.get("full-path")
    if not opf_path:
        return []
    try:
        opf = ET.parse(io.BytesIO(_read_zip_entry(zf, opf_path)))
    except (KeyError, ET.ParseError):
        return []

    base_dir = opf_path.rsplit("/", 1)[0] if "/" in opf_path else ""
    manifest: dict[str, str] = {}
    for item in opf.findall(".//{http://www.idpf.org/2007/opf}item"):
        item_id = item.get("id")
        href = item.get("href")
        if item_id and href:
            manifest[item_id] = f"{base_dir}/{href}" if base_dir else href
    order: list[str] = []
    for ref in opf.findall(".//{http://www.idpf.org/2007/opf}itemref"):
        idref = ref.get("idref")
        if idref and idref in manifest:
            order.append(manifest[idref])
    return order


@register_document_parser(".epub")
def _parse_epub(path: Path, options: dict[str, Any]) -> DocumentIR:
    _require_file(path)
    try:
        zf = zipfile.ZipFile(path)
    except zipfile.BadZipFile as exc:
        raise ValueError(f"document_ir: 无效的 EPUB 文件: {path} ({exc})") from exc

    with zf:
        content_files = [
            name
            for name in zf.namelist()
            if name.lower().endswith((".xhtml", ".html", ".htm"))
        ]
        spine = _epub_spine_order(zf)
        ordered = [name for name in spine if name in content_files]
        if not ordered:
            ordered = sorted(content_files)

        paragraphs: list[str] = []
        title = path.stem
        for name in ordered:
            raw = _read_zip_entry(zf, name)
            text = raw.decode("utf-8", errors="replace")
            document = parse_html(text)
            for node in select_nodes(document, _CONTENT_SELECTOR):
                item = node_text(node)
                if item and item not in paragraphs:
                    paragraphs.append(item)

    title = paragraphs[0] if paragraphs else path.stem
    if title == paragraphs[0] and paragraphs:
        paragraphs = paragraphs[1:]

    return DocumentIR(
        source=path,
        kind=".epub",
        title=title,
        paragraphs=paragraphs,
        metadata={"content_files": len(ordered), "paragraph_count": len(paragraphs)},
    )


__all__ = ["_epub_spine_order"]
