"""S2：document_ir 富文档解析 —— .docx / .pptx / .odt / .epub。

依赖缺失时（python-docx / python-pptx）用 pytest.importorskip 跳过对应用例，
odt/epub 仅用标准库，始终可测。
"""

from __future__ import annotations

import zipfile
from pathlib import Path

import pytest

from omnicrawl.document_ir import DOCUMENT_PARSERS, DocumentIR, parse_document, sniff_document_format


# ── 注册表 ───────────────────────────────────────────────
def test_registry_has_office_parsers() -> None:
    for ext in (".docx", ".pptx", ".odt", ".epub"):
        assert ext in DOCUMENT_PARSERS


def test_sniff_office_formats() -> None:
    assert sniff_document_format(Path("a.docx")) == ".docx"
    assert sniff_document_format(Path("a.pptx")) == ".pptx"
    assert sniff_document_format(Path("a.odt")) == ".odt"
    assert sniff_document_format(Path("a.epub")) == ".epub"


# ── .docx ────────────────────────────────────────────────
def _make_docx(path: Path) -> None:
    import docx

    document = docx.Document()
    document.add_heading("周报标题", level=1)
    document.add_paragraph("本周完成三件事。")
    document.add_paragraph("下周计划两件事。")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "指标"
    table.cell(0, 1).text = "数值"
    table.cell(1, 0).text = "收入"
    table.cell(1, 1).text = "100"
    document.save(str(path))


def test_parse_docx(tmp_path: Path) -> None:
    pytest.importorskip("docx")
    p = tmp_path / "report.docx"
    _make_docx(p)
    doc = parse_document(p)
    assert isinstance(doc, DocumentIR)
    assert doc.kind == ".docx"
    assert doc.title == "周报标题"
    assert "本周完成三件事。" in doc.paragraphs
    assert doc.tables == [[["指标", "数值"], ["收入", "100"]]]
    assert doc.metadata["paragraph_count"] == 2


# ── .pptx ────────────────────────────────────────────────
def _make_pptx(path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "产品发布"
    body = slide.placeholders[1]
    body.text = "要点一\n要点二"
    prs.save(str(path))


def test_parse_pptx(tmp_path: Path) -> None:
    pytest.importorskip("pptx")
    p = tmp_path / "deck.pptx"
    _make_pptx(p)
    doc = parse_document(p)
    assert doc.kind == ".pptx"
    assert doc.title == "产品发布"
    assert "要点一" in doc.paragraphs
    assert "要点二" in doc.paragraphs
    assert doc.metadata["slide_count"] == 1


# ── .odt ─────────────────────────────────────────────────
def _make_odt(path: Path) -> None:
    content = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<office:document-content '
        'xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" '
        'xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
        '<office:body><office:text>'
        "<text:p>ODT 标题</text:p>"
        "<text:p>ODT 正文第一段。</text:p>"
        "<text:p>ODT 正文第二段。</text:p>"
        "</office:text></office:body></office:document-content>"
    )
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("mimetype", "application/vnd.oasis.opendocument.text")
        zf.writestr("content.xml", content)


def test_parse_odt(tmp_path: Path) -> None:
    p = tmp_path / "doc.odt"
    _make_odt(p)
    doc = parse_document(p)
    assert doc.kind == ".odt"
    assert doc.title == "ODT 标题"
    assert "ODT 正文第一段。" in doc.paragraphs
    assert "ODT 正文第二段。" in doc.paragraphs


# ── .epub ────────────────────────────────────────────────
def _make_epub(path: Path) -> None:
    chapter2 = (
        "<html xmlns='http://www.w3.org/1999/xhtml'><body>"
        "<h1>第二章</h1><p>第二章正文。</p></body></html>"
    )
    chapter1 = (
        "<html xmlns='http://www.w3.org/1999/xhtml'><body>"
        "<h1>第一章</h1><p>第一章正文。</p></body></html>"
    )
    container = (
        '<?xml version="1.0"?>'
        '<container xmlns="urn:oasis:names:tc:opendocument:xmlns:container" version="1.0">'
        '<rootfiles><rootfile full-path="OEBPS/content.opf" '
        'media-type="application/oebps-package+xml"/></rootfiles></container>'
    )
    opf = (
        '<?xml version="1.0"?>'
        '<package xmlns="http://www.idpf.org/2007/opf" version="2.0">'
        '<manifest>'
        '<item id="c1" href="chapter1.xhtml" media-type="application/xhtml+xml"/>'
        '<item id="c2" href="chapter2.xhtml" media-type="application/xhtml+xml"/>'
        "</manifest>"
        "<spine><itemref idref=\"c2\"/><itemref idref=\"c1\"/></spine>"
        "</package>"
    )
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("mimetype", "application/epub+zip")
        zf.writestr("META-INF/container.xml", container)
        zf.writestr("OEBPS/content.opf", opf)
        zf.writestr("OEBPS/chapter1.xhtml", chapter1)
        zf.writestr("OEBPS/chapter2.xhtml", chapter2)


def test_parse_epub_respects_spine(tmp_path: Path) -> None:
    p = tmp_path / "book.epub"
    _make_epub(p)
    doc = parse_document(p)
    assert doc.kind == ".epub"
    assert doc.title == "第二章"  # spine 里 c2 在前
    assert doc.paragraphs[0] == "第二章正文。"
    assert "第一章正文。" in doc.paragraphs


def test_parse_epub_bad_zip(tmp_path: Path) -> None:
    p = tmp_path / "bad.epub"
    p.write_bytes(b"not a zip file")
    with pytest.raises(ValueError, match="无效的 EPUB"):
        parse_document(p)


# ── 导出视图 ─────────────────────────────────────────────
def test_to_markdown_docx(tmp_path: Path) -> None:
    pytest.importorskip("docx")
    p = tmp_path / "r.docx"
    _make_docx(p)
    doc = parse_document(p)
    md = doc.to_markdown()
    assert md.startswith("# 周报标题")
    assert "| 指标 | 数值 |" in md
